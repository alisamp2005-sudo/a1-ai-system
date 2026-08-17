"""
Document Processor — извлечение текста из файлов различных форматов.

Поддерживаемые форматы:
- PDF (.pdf)
- Word (.docx, .doc)
- Excel (.xlsx, .xls)
- Текст (.txt, .md, .csv, .json, .xml)
- Изображения (.jpg, .jpeg, .png, .tiff, .bmp) — через OCR
- PowerPoint (.pptx)
- RTF (.rtf)
"""

import logging
import os
import tempfile
from pathlib import Path
from typing import Optional, Tuple

logger = logging.getLogger(__name__)

# Поддерживаемые форматы
SUPPORTED_FORMATS = {
    # Документы
    ".pdf": "PDF документ",
    ".docx": "Word документ",
    ".doc": "Word документ (старый формат)",
    ".xlsx": "Excel таблица",
    ".xls": "Excel таблица (старый формат)",
    ".pptx": "PowerPoint презентация",
    # Текстовые
    ".txt": "Текстовый файл",
    ".md": "Markdown файл",
    ".csv": "CSV таблица",
    ".json": "JSON файл",
    ".xml": "XML файл",
    ".rtf": "RTF документ",
    ".html": "HTML страница",
    ".htm": "HTML страница",
    # Изображения (OCR)
    ".jpg": "Изображение (OCR)",
    ".jpeg": "Изображение (OCR)",
    ".png": "Изображение (OCR)",
    ".tiff": "Изображение (OCR)",
    ".tif": "Изображение (OCR)",
    ".bmp": "Изображение (OCR)",
    ".webp": "Изображение (OCR)",
}


def get_supported_extensions() -> list:
    """Return list of supported file extensions."""
    return list(SUPPORTED_FORMATS.keys())


def is_supported(filename: str) -> bool:
    """Check if file format is supported."""
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_FORMATS


def get_format_description(filename: str) -> str:
    """Get human-readable format description."""
    ext = Path(filename).suffix.lower()
    return SUPPORTED_FORMATS.get(ext, "Неизвестный формат")


async def extract_text(file_path: str) -> Tuple[str, str]:
    """
    Extract text from a file.

    Args:
        file_path: Path to the file

    Returns:
        Tuple of (extracted_text, format_description)
    """
    ext = Path(file_path).suffix.lower()
    format_desc = SUPPORTED_FORMATS.get(ext, "Неизвестный")

    try:
        if ext == ".pdf":
            text = await _extract_pdf(file_path)
        elif ext in (".docx",):
            text = await _extract_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            text = await _extract_excel(file_path)
        elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".rtf"):
            text = await _extract_text_file(file_path)
        elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp"):
            text = await _extract_image_ocr(file_path)
        elif ext == ".pptx":
            text = await _extract_pptx(file_path)
        else:
            text = ""

        if not text or len(text.strip()) < 10:
            return "", format_desc

        return text.strip(), format_desc

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return "", format_desc


async def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyPDF2 or pdfplumber."""
    try:
        import PyPDF2
        text_parts = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts)
    except ImportError:
        pass

    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts)
    except ImportError:
        pass

    # Fallback: use subprocess with pdftotext
    import subprocess
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", file_path, "-"],
            capture_output=True, text=True, timeout=30
        )
        return result.stdout
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return ""


async def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX")
        return ""


async def _extract_excel(file_path: str) -> str:
    """Extract text from Excel."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"=== Лист: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        wb.close()
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("openpyxl not installed, cannot extract Excel")
        return ""


async def _extract_text_file(file_path: str) -> str:
    """Extract text from plain text files."""
    encodings = ["utf-8", "cp1251", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


async def _extract_image_ocr(file_path: str) -> str:
    """Extract text from image using Vision service (MLX)."""
    import httpx

    vision_url = os.getenv("VISION_URL", "http://host.docker.internal:11435")

    try:
        import base64
        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                f"{vision_url}/analyze",
                json={
                    "image_base64": image_data,
                    "prompt": (
                        "Extract ALL text from this document image. "
                        "Return only the text content, preserving structure. "
                        "If it's a table, format as rows. "
                        "If no text found, say 'NO_TEXT'."
                    ),
                },
            )

            if response.status_code == 200:
                data = response.json()
                text = data.get("analysis", "")
                if "NO_TEXT" in text:
                    return ""
                return text
    except Exception as e:
        logger.warning(f"OCR via Vision service failed: {e}")

    return ""


async def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                text_parts.append(f"--- Слайд {i} ---")
                text_parts.extend(slide_texts)
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX")
        return ""
